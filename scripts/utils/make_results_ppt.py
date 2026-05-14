"""Generate a PPT summarizing all experiments & results to date.

Run: PYTHONPATH=src /home/abhor/miniconda3/envs/agentic/bin/python \
       scripts/utils/make_results_ppt.py
Output: runs/analysis/experiments_summary.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT = Path("runs/analysis/experiments_summary.pptx")

NAVY = RGBColor(0x0B, 0x2E, 0x4F)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xD6, 0x27, 0x28)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title_slide(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # background bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(2))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sub = s.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(12), Inches(2))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xCF, 0xE2, 0xF3)
    return s


def add_section(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), prs.slide_width, Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return s


def add_content(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # title
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return s


def add_bullets(slide, bullets, left=0.5, top=1.1, width=12.3, height=6.0, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(b, tuple):
            text, level = b
        else:
            text, level = b, 0
        p.text = ("• " if level == 0 else "   – ") + text
        p.level = level
        p.font.size = Pt(size if level == 0 else size - 2)
        p.font.color.rgb = NAVY if level == 0 else GREY
        p.space_after = Pt(4)


def add_table(slide, headers, rows, left=0.5, top=1.2, width=12.3, height=5.5,
              header_fill=NAVY, highlight_rows=None, highlight_color=None):
    highlight_rows = highlight_rows or set()
    highlight_color = highlight_color or RGBColor(0xFF, 0xF2, 0xCC)
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers),
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height)).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = header_fill
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            if i - 1 in highlight_rows:
                c.fill.solid()
                c.fill.fore_color.rgb = highlight_color
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.color.rgb = NAVY
    return tbl


# ---------------------------------------------------------------- 1. Title
add_title_slide(
    "Google Calendar MCP Agent",
    "All experiments & results — 2026-03 to 2026-04   |   Qwen2.5-1.5B → Qwen3-14B",
)

# ---------------------------------------------------------------- 2. Problem
s = add_content("Problem & Approach")
add_bullets(s, [
    "Goal: tool-calling agent for Google Calendar — 7 tools, multi-turn, calendar-aware",
    "Stack: SFT (Unsloth + TRL, LoRA) → optional RL (GRPO via ART 0.5.17 + vLLM) / DPO",
    "Eval: Gemini-2.0-flash-as-judge compares calendar state before/after vs expected behavior",
    "7 task categories: RelTime, IR, Modifier, Schedule, Vague, Chaos, Complex Logic",
    "Datasets:",
    ("SFT: 6,947 augmented trajectories from 114 calendars (gemini-2.5-pro teacher)", 1),
    ("RL:  622 scenarios across 50 calendars (binary correct/incorrect rewards)", 1),
    ("Test (held-out, canonical): 692 queries × 49 fresh calendars — zero overlap", 1),
    "Hardware journey: TITAN X 12 GiB (1.5B era) → Blackwell 96 GiB MIG-partitioned 4×24 GiB",
])

# ---------------------------------------------------------------- 3. Pipeline
s = add_content("Pipeline Overview")
add_bullets(s, [
    "Data generation",
    ("114 calendars + 7-category prompt templates → gemini-2.5-pro trajectories (compact tool-result format)", 1),
    ("Augmentation: entity substitution (2×) + paraphrasing (2-5×) → 6,947 balanced trajectories", 1),
    "SFT (Qwen3-14B + 4-bit bnb + LoRA r=64, bf16, /no_think system prompt)",
    ("5 epochs, cosine LR 2e-4, max_seq_len=4096, loss-masked on assistant + tool_call tokens", 1),
    "RL — GRPO (openpipe-art 0.5.17)",
    ("rollouts_per_group=8, num_generations=4, lr=5e-6, β=0, binary reward from Gemini judge", 1),
    ("Multi-tenant: MIG slice 0/1/2/3, auto_restart wrapper, Patch G deadlock recovery", 1),
    "DPO — pairs mined from RL rollouts (chosen=correct, rejected=incorrect) on same scenario",
    "Eval — vLLM serve (hermes parser) → batched judge across 692 queries × 7 categories",
])

# ---------------------------------------------------------------- 4. Section: 1.5B Era
add_section("Phase 1 — Qwen2.5-1.5B Era  (TITAN X 12 GiB)")

# 5. SFT v3 / v5
s = add_content("SFT v3 → v5 (1.5B)  — overall acc on RL data, 280 queries")
add_table(s, ["Run", "Best ckpt", "RL-data acc", "Notes"], [
    ["SFT v3 (compact)",  "ckpt-933 (ep 6)",   "42.5%",  "Baseline; trained on 1,164 trajectories"],
    ["SFT v5 (augmented)", "ckpt-6152 (ep 4)", "74.6%",  "Augmented to 6,947; narration-merge fix"],
], top=1.3, height=1.2)
tb = s.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.3), Inches(4.0))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate([
    "v5 fixes that mattered:",
    "   – Phase-4 augmentation bug: 29.4% of training data was corrupted (get_current_time results were "
    "swapped with random event names) → first run scored 18-21%; fix lifted +50pp",
    "   – Narration-merge in trajectory_to_messages — assistant narration merged into tool_call message",
    "   – Cosine-with-restarts produced even/odd LR oscillation; epoch 4,6,8 > 3,5,7,9",
    "Best v5 per-category at epoch 4: IR 95%, RelTime 92.5%, Schedule 87.5%; Complex 32.5%, Vague 62.5%",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(15)
    p.font.bold = (i == 0)
    p.font.color.rgb = NAVY

# 6. RL on 1.5B
s = add_content("RL Runs 1-3 on 1.5B  (GRPO, binary reward, β=0)")
add_table(s, ["Run", "Target cat.", "Steps", "Overall (RL data)", "Δ vs prior", "Verdict"], [
    ["RL1 Modifier",       "Modifier",   "395 (5 ep)", "42.5% → 47.9%", "+5.4 pp", "Target +12.5pp"],
    ["RL2 Information Retr","IR",         "234",        "47.9% → 47.5%", "−0.4 pp", "Flat overall"],
    ["RL3 Vague",          "Vague",      "440 (5 ep)", "Flat 66-69%",   "—",        "No improvement"],
    ["SFT-on-RL recovery", "All",        "1 epoch",    "47.9% → 39.3%", "−8.6 pp", "Overwrote RL gains"],
], top=1.3, height=2.3)
add_bullets(s, [
    "Failure mode (RL3 trajectory analysis, 2,464 rollouts, 806 incorrect):",
    ("79% incomplete/wrong answer — right tools, wrong interpretation", 1),
    ("16% action instead of query (created/updated when should have read)", 1),
    ("→ Comprehension ceiling, not procedure. 1.5B can't replicate teacher's semantic reasoning.", 1),
    "Other 1.5B failure modes: bimodal skip rate (50-65%), catastrophic forgetting (Vague 50→35%)",
], top=4.0, size=15)

# ---------------------------------------------------------------- 7. Section: 14B
add_section("Phase 2 — Switch to Qwen3-14B  (Blackwell, 4×24 GiB MIG)")

# 8. Stack & key changes
s = add_content("Stack upgrade & key config changes  (2026-04-15)")
add_table(s, ["Component", "Old (1.5B / TITAN X)", "New (14B / Blackwell)"], [
    ["Model",            "Qwen2.5-1.5B",            "Qwen3-14B + /no_think"],
    ["Quantization",     "fp16",                     "4-bit bnb + bf16 + FA2"],
    ["LoRA rank",        "8 (RL OOM @ 64)",          "64 (SFT) / 64 (RL with sleep_mode)"],
    ["max_model_len",    "3076",                     "4096"],
    ["num_generations",  "2",                        "4"],
    ["torch / vLLM",     "2.x / 0.7.3 (sm_61 build)", "2.10.0+cu128 / 0.19.0"],
    ["unsloth / TRL",    "older",                    "2026.4.4 / 0.24.0"],
    ["openpipe-art",     "0.5.4",                    "0.5.17 (with patches D, E, G, H, I)"],
], top=1.2, height=4.0)
add_bullets(s, [
    "vLLM tool-call parser MUST be hermes (not qwen3_xml) — Qwen3 outputs <tool_call>…</tool_call>",
    "RL adapters must be served via --enable-lora; merging RL LoRA → fp16 produces 0% acc",
], top=5.5, size=14)

# 9. SFT v6 full results
s = add_content("SFT v6 (Qwen3-14B) — full per-checkpoint eval on held-out test_data (692 queries)")
add_table(s, ["Epoch", "Ckpt", "Overall", "Complex", "Chaos", "IR", "Modif", "RelTime", "Schedule", "Vague"], [
    ["1", "1553", "72.4%", "44%", "67%", "88%", "70%", "88%", "78%", "70%"],
    ["2", "3106", "78.3%", "54%", "77%", "89%", "87%", "93%", "68%", "80%"],
    ["3", "4659", "80.1%", "55%", "77%", "92%", "85%", "93%", "77%", "82%"],
    ["4", "6212", "78.9%", "59%", "79%", "89%", "89%", "90%", "67%", "79%"],
    ["5", "7765", "79.2%", "52%", "83%", "90%", "87%", "92%", "70%", "80%"],
], top=1.2, height=2.8, highlight_rows={2})
add_bullets(s, [
    "BEST: ckpt-4659 (epoch 3) at 80.1% on held-out test_data — current production ckpt",
    "Cross-set discrepancy: ckpt-6212 wins on RL-data benchmark (82.5%); ckpt-4659 wins on test_data",
    "Eval loss is misleading for ckpt selection — track full eval, not loss minimum",
    "Versus 1.5B SFT v5 best (74.6% RL-data): +7.9 pp overall after switch to 14B",
], top=4.2, size=15)

# 10. SFT v6 vs Instruct baseline per-category
s = add_content("SFT v6 vs Qwen3-14B-Instruct baseline — per-category gains")
add_table(s, ["Category", "Instruct base", "SFT v6 best", "Δ", "Comment"], [
    ["RelTime",  "95%", "93%",  "−2 pp",  "Saturated by base model"],
    ["IR",       "60%", "92%",  "+32 pp", "Largest absolute lift"],
    ["Modifier", "70%", "85%",  "+14 pp", "Strong"],
    ["Chaos",    "30%", "77%",  "+47 pp", "Calendar-specific knowledge"],
    ["Vague",    "70%", "82%",  "+11 pp", "Solid"],
    ["Schedule", "72%", "77%",  "+5 pp",  "Marginal"],
    ["Complex",  "41%", "55%",  "+14 pp", "Weakest absolute (multi-step reasoning)"],
    ["Overall",  "63.0%","80.1%","+17.1pp", "—"],
], top=1.2, height=4.0, highlight_rows={7})
add_bullets(s, [
    "SFT lifts the hardest, knowledge-bound categories most (Chaos +47, IR +32)",
    "Categories already near ceiling (RelTime) cannot be moved further by domain SFT",
], top=5.6, size=14)

# 11. DPO experiments
s = add_content("DPO experiments (paused 2026-04-25)")
add_table(s, ["Model", "Best epoch", "Test-set acc", "vs starting point", "Verdict"], [
    ["DPO-from-SFT v6",     "ep 1 / ep 2", "79.3%", "−0.8 pp vs SFT 80.1%",  "Within noise; no convincing win"],
    ["DPO-from-Instruct",   "ep 2",        "61.3%", "−1.7 pp vs Instruct 63%", "Actively regresses (likelihood displacement)"],
], top=1.3, height=1.6)
add_bullets(s, [
    "Pair source: 1,913 mined pairs from 5,506 ART rollouts on RL-data scenarios",
    "Why DPO-from-Instruct failed (predicted ex-ante):",
    ("π_ref assigns near-zero prob to <tool_call> XML for un-SFT'd Instruct → loss ratio dominated by", 1),
    ("likelihood displacement (Pal et al. 2024). Chaos 30→26%, Vague 70→58% confirms.", 1),
    "Why DPO-from-SFT washed out:",
    ("Pair saturation — chosen trajectories already near SFT v6's ceiling", 1),
    ("Per-category: Schedule +9 pp gain offset by Complex −7 pp regression; net ~ 0", 1),
    "Decision: pause DPO; explore RFT / iterative-DPO / Complex-targeted SFT instead",
], top=3.2, size=14)

# 12. RL on 14B
s = add_content("RL on Qwen3-14B  (GRPO, ART 0.5.17)")
add_bullets(s, [
    "Setup: all 7 categories joint, 622 scenarios, rollouts_per_group=8, num_gens=4, lr=5e-6, β=0",
    "Plan: 622 × 20 epochs = 12,440 steps. Currently paused at step 9220.",
    "Mid-run (step 2325, 2026-04-17 snapshot): mean reward 0.70, recent-500 ≈ 0.76",
    ("Per-category recent: RelTime 93%, Schedule 75%, IR 74%, Vague 71%, Modifier 68%, Chaos 61%, Complex 47%", 1),
    "Incidents:",
    ("2026-04-17 — 57h hang inside model.train() at step 2325; root cause unresolved, restarted 2026-04-20", 1),
    ("2026-04-25 — reward cliff; led to multi-tenant hardening + Patch K v2 (milestone preservation)", 1),
    "Currently paused: real-RL adapter at .art/calendar-agent/.../checkpoints/9220/",
    "Active follow-ups: rl_adaptive variant, rl_grpo from base + from sft4659 (parallel 2026-04-26 runs)",
], size=15)

# 13. ART patches summary
s = add_content("ART 0.5.17 runtime patches  (src/calendar_agent/art_patches.py)")
add_table(s, ["Patch", "Target", "Purpose"], [
    ["D", "_calculate_logprobs",      "entropy detach from autograd; remove chunk_size assertion"],
    ["E", "_prepare_backend_for_training", "guard done_callback against error/cancel (prevents OOM on health-check timeout)"],
    ["G", "inputs_queue.get",         "300s timeout + exit(42) on asyncio queue deadlock; auto_restart relaunches"],
    ["H", "tokenize_trajectory",      "REPAIR Qwen3 empty <think></think>-only finals → patch content=' ', retain for negative gradient"],
    ["I", "(asyncio nest)",           "nest_asyncio shim — replaces older A/B/C patch trio"],
    ["F", "(optional)",               "LoRA injection helper for resuming RL from arbitrary adapter"],
], top=1.2, height=3.5)
add_bullets(s, [
    "Patches A, B, C removed (fixed upstream / replaced by nest_asyncio approach)",
    "Patch H upgrade (2026-04-24): drop-only → repair-then-train. Was losing ~7% of rollouts (~1 per 15 steps).",
    "Import calendar_agent.art_patches BEFORE import art in rl_train.py",
], top=4.9, size=14)

# 14. Headline test-set comparison
s = add_content("Headline — Held-out test_data benchmark  (49 cals × 692 queries)")
add_table(s, ["Model", "Best", "Accuracy", "Notes"], [
    ["Qwen3-14B-Instruct (no calendar training)", "—",       "63.0%",                "Format-strong, knowledge-weak"],
    ["DPO-from-Instruct",                          "ep 2",   "61.3% (−1.7 pp)",      "Likelihood displacement"],
    ["DPO-from-SFT",                               "ep 1/2", "79.3% (−0.8 pp)",      "Within noise of SFT"],
    ["SFT v6 ckpt-4659 (epoch 3)",                 "—",      "80.1%   ← BEST",        "Production checkpoint"],
    ["SFT v6 ckpt-6212 (epoch 4)",                 "—",      "78.9%",                 "Best on RL-data set (82.5%)"],
], top=1.3, height=2.8, highlight_rows={3})
add_bullets(s, [
    "test_data created 2026-04-24 because rl_data was contaminated for DPO eval (74% overlap with mined pairs)",
    "Canonical benchmark going forward; old rl_data numbers retained for SFT-vs-SFT historical comparison",
    "SE ≈ 1.5 pp on 692 queries — differences below ~3 pp are within noise",
], top=4.4, size=14)

# 15. Per-category status
s = add_content("Per-category status  (SFT v6 ckpt-4659 on test_data)")
add_table(s, ["Category", "Acc", "Status", "Notes"], [
    ["RelTime",  "95%", "saturated", "Improved by base model alone"],
    ["IR",       "92%", "saturated", "Largest SFT lift over baseline"],
    ["Modifier", "85%", "saturated", "Strong"],
    ["Chaos",    "77%", "strong",    "Edge-case calendars; 47-pp SFT lift"],
    ["Vague",    "82%", "room",      "Contextual / underspecified queries"],
    ["Schedule", "77%", "room",      "Event creation format precision"],
    ["Complex",  "55%", "weakest",   "Multi-step reasoning + tool chaining"],
], top=1.3, height=3.3)
add_bullets(s, [
    "Top improvement target: Complex Logic (multi-step reasoning + conflict resolution)",
    "Secondary: Schedule + Vague (room above 75-80% with focused data)",
], top=5.0, size=15)

# 16. Open threads / next
s = add_content("Open threads & next experiments")
add_bullets(s, [
    "Local judge (Qwen3-7B SFT) — replaces ~99k Gemini API calls per RL run; baseline eval in progress",
    "ART asyncio deadlock — workaround deployed (Patch G/I); upstream issue not yet filed",
    "RL beyond GRPO+binary rewards:",
    ("RFT / Expert iteration — filter correct rollouts → SFT (avoids DPO ratio failure modes)", 1),
    ("Iterative DPO — regenerate pairs from current SFT v6 ckpt-4659 (on-policy)", 1),
    ("Dr. GRPO / non-zero β — stop catastrophic forgetting; better advantage for binary reward", 1),
    ("Complex-category focused SFT — directly attack the 55% bottleneck", 1),
    "Multi-tenant hardening (2026-04-26): centralized auto_restart, slice_map, telemetry module, /rl-status & /rl-stop skills",
], size=16)

# 17. Summary
s = add_content("Summary")
add_bullets(s, [
    "End-to-end pipeline (data gen → SFT → RL/DPO → judge eval) running on 4× MIG Blackwell partitions",
    "1.5B era hit a comprehension ceiling at 47.9% (RL data); 14B switch lifted SFT to 80.1% on held-out",
    "Current best model: SFT v6 Qwen3-14B ckpt-4659 (epoch 3) — 80.1% on 692-query held-out benchmark",
    "DPO experiments paused: neither variant convincingly beats its starting point",
    "RL on 14B paused at step 9220 after reward cliff; multi-tenant safeguards now in place for resume",
    "Next priorities: Local judge, RFT / iterative DPO, Complex-category targeted SFT",
], size=20)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT}  ({OUT.stat().st_size/1024:.1f} KiB, {len(prs.slides)} slides)")
